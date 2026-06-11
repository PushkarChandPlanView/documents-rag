"""
format_editors.py
Apply natural-language edits to the original source file bytes.
Returns updated bytes for supported formats, or None for unsupported ones
(PDF / images) so the caller can fall back to text-only editing.
"""
from __future__ import annotations

import difflib
import io
from typing import Optional


def _build_replacement_map(original_text: str, proposed_text: str) -> dict[str, str]:
    """
    Build a line-level replacement map from original_text → proposed_text
    using difflib.SequenceMatcher so we only touch lines that actually changed.
    """
    orig_lines = original_text.splitlines()
    new_lines = proposed_text.splitlines()
    replacements: dict[str, str] = {}

    matcher = difflib.SequenceMatcher(None, orig_lines, new_lines, autojunk=False)
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == "replace":
            # Pair up changed lines 1-to-1 as far as possible
            for orig_line, new_line in zip(orig_lines[i1:i2], new_lines[j1:j2]):
                if orig_line.strip():
                    replacements[orig_line] = new_line
        elif tag == "delete":
            for orig_line in orig_lines[i1:i2]:
                if orig_line.strip():
                    replacements[orig_line] = ""

    return replacements


# ── plain text / markdown / csv ───────────────────────────────────────────────

def _edit_text(proposed_text: str) -> bytes:
    return proposed_text.encode("utf-8")


# ── DOCX ─────────────────────────────────────────────────────────────────────

def _apply_replacements_to_runs(paragraph, replacements: dict[str, str]) -> None:
    """
    Apply replacements to a docx paragraph while preserving per-run formatting.
    We reconstruct each run's text by replacing matching substrings.
    """
    for run in paragraph.runs:
        text = run.text
        for orig, new in replacements.items():
            if orig and orig in text:
                text = text.replace(orig, new)
        run.text = text


def _edit_docx(original_bytes: bytes, replacements: dict[str, str]) -> bytes:
    from docx import Document  # type: ignore

    doc = Document(io.BytesIO(original_bytes))

    # Body paragraphs
    for para in doc.paragraphs:
        _apply_replacements_to_runs(para, replacements)

    # Table cells
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _apply_replacements_to_runs(para, replacements)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _edit_xlsx(original_bytes: bytes, proposed_text: str) -> bytes:
    """
    Handles two cases:
    1. proposed_text is a JSON array of sheet specs → full structural rebuild
       (supports add/modify sheets, add columns, add rows)
    2. proposed_text is plain text → fall back to line-level replacement
    """
    import json
    from openpyxl import load_workbook, Workbook  # type: ignore

    # ── Try JSON structural mode ──────────────────────────────────────────────
    json_text = proposed_text.strip()
    # Strip markdown code fences if LLM wrapped the JSON
    if json_text.startswith("```"):
        lines = json_text.splitlines()
        json_text = "\n".join(l for l in lines if not l.startswith("```"))

    try:
        sheets_spec = json.loads(json_text)
        if isinstance(sheets_spec, list) and all(isinstance(s, dict) for s in sheets_spec):
            wb = Workbook()
            wb.remove(wb.active)  # remove default empty sheet
            for spec in sheets_spec:
                ws = wb.create_sheet(title=str(spec.get("name", "Sheet")))
                columns = spec.get("columns") or []
                if columns:
                    ws.append(columns)
                for row in (spec.get("rows") or []):
                    ws.append([None if v is None else v for v in row])
            buf = io.BytesIO()
            wb.save(buf)
            return buf.getvalue()
    except (json.JSONDecodeError, Exception):
        pass  # fall through to text-replacement mode

    # ── Fallback: line-level text replacement ─────────────────────────────────
    from openpyxl import load_workbook  # type: ignore

    orig_text = proposed_text  # in fallback, proposed_text IS the replacement target
    wb = load_workbook(io.BytesIO(original_bytes))
    replacements = _build_replacement_map(orig_text, proposed_text)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    new_val = cell.value
                    for orig, new in replacements.items():
                        if orig and orig in new_val:
                            new_val = new_val.replace(orig, new)
                    cell.value = new_val
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _edit_pptx(original_bytes: bytes, replacements: dict[str, str]) -> bytes:
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(original_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text = run.text
                    for orig, new in replacements.items():
                        if orig and orig in text:
                            text = text.replace(orig, new)
                    run.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Public entry-point ────────────────────────────────────────────────────────

PLAIN_TEXT_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
}

UNSUPPORTED_TYPES_PREFIX = (
    "application/pdf",
    "image/",
)

MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def apply_edit_to_file(
    original_bytes: bytes,
    mime_type: str,
    original_text: str,
    proposed_text: str,
) -> Optional[bytes]:
    """
    Returns updated file bytes for supported MIME types, or None when the
    format cannot be modified (PDF / images) and the caller should fall back
    to text-only editing.
    """
    mt = (mime_type or "").lower().strip()

    # Unsupported — signal fallback
    if mt == "application/pdf" or mt.startswith("image/"):
        return None

    # Plain text variants — just return the proposed text as UTF-8
    if mt in PLAIN_TEXT_TYPES:
        return _edit_text(proposed_text)

    if mt == MIME_XLSX:
        return _edit_xlsx(original_bytes, proposed_text)

    replacements = _build_replacement_map(original_text, proposed_text)

    if mt == MIME_DOCX:
        return _edit_docx(original_bytes, replacements)

    if mt == MIME_PPTX:
        return _edit_pptx(original_bytes, replacements)

    # Unknown binary format — fall back to text-only
    return None
