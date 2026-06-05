import io
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Pt


@dataclass
class ExtractionResult:
    text: str
    page_count: int


def extract(file_bytes: bytes) -> ExtractionResult:
    prs = Presentation(io.BytesIO(file_bytes))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    lines.append(text)
        if lines:
            slides.append(f"[Slide {i}]\n" + "\n".join(lines))
    full_text = "\n\n".join(slides)
    return ExtractionResult(text=full_text, page_count=len(prs.slides))
